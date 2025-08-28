# modules/loader.py
# -*- coding: utf-8 -*-

import os
from typing import List, Optional
from langchain.schema import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

# PDF 轻量解析：PyMuPDF
from langchain_community.document_loaders import PyMuPDFLoader

# 轻量依赖（不用 unstructured 也能跑）
import docx2txt
from pptx import Presentation
from bs4 import BeautifulSoup
from unstructured.partition.auto import partition

# ------------------------------------------------------------
# 切换策略与支持格式
# ------------------------------------------------------------
SUPPORTED_EXTS = {".txt", ".md", ".html", ".htm", ".docx", ".pptx", ".pdf"}

# 采用轻量文件解析


# ------------------------------------------------------------
# 工具函数
# ------------------------------------------------------------
def _mk_docs(texts: List[str], path: str) -> List[Document]:
    """把多段纯文本封装为 Document 列表，并补充统一 metadata。"""
    basename = os.path.basename(path)
    apath = os.path.abspath(path)
    ext = os.path.splitext(path)[-1].lower()
    docs: List[Document] = []
    for t in texts:
        t = (t or "").strip()
        if not t:
            continue
        docs.append(
            Document(
                page_content=t,
                metadata={"source": basename, "path": apath, "ext": ext},
            )
        )
    return docs


# ------------------------------------------------------------
# 轻量解析实现（无 unstructured）
# ------------------------------------------------------------
def _load_txt(path: str) -> List[Document]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        txt = f.read()
    return _mk_docs([txt], path)


def _load_md(path: str) -> List[Document]:
    # 以原始 Markdown 文本参与检索；如需转纯文本，可进一步处理
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        md = f.read()
    return _mk_docs([md], path)


def _load_html(path: str) -> List[Document]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    soup = BeautifulSoup(html, "html.parser")
    text = soup.get_text(separator="\n")
    return _mk_docs([text], path)


def _load_docx(path: str) -> List[Document]:
    text = docx2txt.process(path) or ""
    return _mk_docs([text], path)


def _load_pptx(path: str) -> List[Document]:
    prs = Presentation(path)
    slide_texts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        buf = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                buf.append(shape.text)
        if buf:
            slide_texts.append(f"[Slide {i}]\n" + "\n".join(buf))
    return _mk_docs(slide_texts or [""], path)


def _load_pdf(path: str) -> List[Document]:
    loader = PyMuPDFLoader(path)
    docs = loader.load()
    # 统一补 metadata
    basename = os.path.basename(path)
    apath = os.path.abspath(path)
    for d in docs:
        d.metadata.setdefault("source", basename)
        d.metadata.setdefault("path", apath)
        d.metadata.setdefault("ext", ".pdf")
    return docs


def _load_light(path: str) -> List[Document]:
    ext = os.path.splitext(path)[-1].lower()
    if ext == ".txt":
        return _load_txt(path)
    if ext == ".md":
        return _load_md(path)
    if ext in (".html", ".htm"):
        return _load_html(path)
    if ext == ".docx":
        return _load_docx(path)
    if ext == ".pptx":
        return _load_pptx(path)
    if ext == ".pdf":
        return _load_pdf(path)
    raise ValueError(f"暂不支持的文件类型: {ext}")




# ------------------------------------------------------------
# 对外 API
# ------------------------------------------------------------
def load_multiple_documents(
    paths: List[str],
    *,
    limit: int = 5,
    prefer_mode: str = "elements",
    prefer_strategy: str = "fast",
    infer_table_structure: bool = False,
) -> List[Document]:
    """
    批量加载文档，支持多种格式（txt, md, html, docx, pptx, pdf）。
    """
    if len(paths) > limit:
        raise ValueError(f"最多仅支持上传 {limit} 个文件，请减少文件数量后重试。")

    all_docs: List[Document] = []
    for path in paths:
        ext = os.path.splitext(path)[-1].lower()
        if ext not in SUPPORTED_EXTS:
            raise ValueError(f"暂不支持的文件类型: {ext}（支持：{', '.join(sorted(SUPPORTED_EXTS))}）")

        docs: Optional[List[Document]] = None

        docs = _load_light(path)
        if not docs:
            continue
        
        all_docs.extend(docs or [])

    return all_docs


def split_documents(documents: List[Document], chunk_size=500, chunk_overlap=50):
    """将加载后的文档切分为小段落块"""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
    )
    return splitter.split_documents(documents)
